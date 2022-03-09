import sys

from pptx import Presentation


for slide in Presentation(sys.argv[1]).slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print(run.text)